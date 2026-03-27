"""Extract structured data from assignment deliverables."""
from __future__ import annotations

import json
import logging
import os
import re
import zipfile
from collections import Counter
from dataclasses import dataclass, field
from typing import Any, Dict, List, Optional, Sequence
import xml.etree.ElementTree as ET

try:  # pragma: no cover - dependency resolved at runtime
    import docx  # type: ignore[import]
except ImportError:  # pragma: no cover
    docx = None  # type: ignore[assignment]

try:  # pragma: no cover
    import pandas as pd  # type: ignore[import]
except ImportError:  # pragma: no cover
    pd = None  # type: ignore[assignment]

try:  # pragma: no cover
    import pptx  # type: ignore[import]
except ImportError:  # pragma: no cover
    pptx = None  # type: ignore[assignment]

try:  # pragma: no cover - prance for OpenAPI parsing
    from prance import ResolvingParser  # type: ignore[import]
except ImportError:  # pragma: no cover
    ResolvingParser = None  # type: ignore[assignment]


from .exceptions import MappingDataError


def _ensure_docx() -> None:
    global docx
    if docx is not None:
        return
    try:
        import importlib

        docx = importlib.import_module("docx")  # type: ignore[assignment]
    except ImportError as exc:  # pragma: no cover
        raise MappingDataError("python-docx is not installed; cannot parse Word documents.") from exc


def _ensure_pptx() -> None:
    global pptx
    if pptx is not None:
        return
    try:
        import importlib

        pptx = importlib.import_module("pptx")  # type: ignore[assignment]
    except ImportError as exc:  # pragma: no cover
        raise MappingDataError("python-pptx is not installed; cannot parse presentations.") from exc


def _ensure_prance() -> None:
    global ResolvingParser
    if ResolvingParser is not None:
        return
    try:
        import importlib

        prance_module = importlib.import_module("prance")  # type: ignore[assignment]
        ResolvingParser = prance_module.ResolvingParser  # type: ignore[assignment]
    except ImportError as exc:  # pragma: no cover
        raise MappingDataError("prance is not installed; cannot parse Swagger/OpenAPI specs.") from exc

logger = logging.getLogger("agentic_orchestrator_auto.mapping.parsers")
_ANGLE_PLACEHOLDER_PATTERN = re.compile(r"<[^<>]+>")
_BRACKET_PLACEHOLDER_PATTERN = re.compile(r"\[[^\[\]\r\n]+?\]")
_WORD_NAMESPACE = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
_ENDNOTE_PATH = "word/endnotes.xml"

_FIELD_NAME_HINTS: Sequence[str] = ("name", "field", "column", "attribute", "label", "title")
_FIELD_PATH_HINTS: Sequence[str] = ("jsonpath", "path", "pointer", "selector", "key", "route", "xpath")
_FIELD_DESCRIPTION_HINTS: Sequence[str] = ("description", "desc", "details", "definition", "summary", "meaning")
_FIELD_SAMPLE_HINTS: Sequence[str] = ("sample", "example", "value", "default", "format")
_EXCLUDE_METADATA_HINTS: Sequence[str] = tuple(
    set(_FIELD_NAME_HINTS)
    | set(_FIELD_PATH_HINTS)
    | set(_FIELD_DESCRIPTION_HINTS)
    | set(_FIELD_SAMPLE_HINTS)
)
_DESCRIPTOR_LIMIT_PER_SHEET = 400
_SUPERSCRIPT_MAP = {
    "\u2070": "0",
    "\u00b9": "1",
    "\u00b2": "2",
    "\u00b3": "3",
    "\u2074": "4",
    "\u2075": "5",
    "\u2076": "6",
    "\u2077": "7",
    "\u2078": "8",
    "\u2079": "9",
}
_SUPERSCRIPT_CLASS = "".join(_SUPERSCRIPT_MAP.keys())
_SUPERSCRIPT_CHARS = set(_SUPERSCRIPT_CLASS)
_SUPERSCRIPT_DEF_PATTERN = re.compile(rf"^(?P<sup>[{_SUPERSCRIPT_CLASS}])[\s\-.):]*?(?P<body>.+)")
_SUPERSCRIPT_DIGIT_PATTERN = re.compile(r"\d+")


def _clean_text(value: Any) -> str:
    if value is None:
        return ""
    try:
        if pd is not None and pd.isna(value):  # type: ignore[attr-defined]
            return ""
    except Exception:  # pragma: no cover
        pass
    text = str(value).strip()
    if not text or text.lower() in {"nan", "none", "null"}:
        return ""
    return text


def _extract_value(mapping: Dict[str, Any], hints: Sequence[str]) -> Optional[str]:
    for key, raw_value in mapping.items():
        lowered = str(key).lower()
        if any(hint in lowered for hint in hints):
            text = _clean_text(raw_value)
            if text:
                return text
    return None


def _filter_metadata(mapping: Dict[str, Any]) -> Dict[str, Any]:
    filtered: Dict[str, Any] = {}
    for key, value in mapping.items():
        lowered = str(key).lower()
        if any(token in lowered for token in _EXCLUDE_METADATA_HINTS):
            continue
        filtered[str(key)] = value
    return filtered


@dataclass
class ExcelObjectDescriptor:
    sheet: str
    row_index: int
    column: Optional[str]
    name: Optional[str]
    description: Optional[str]
    path: Optional[str]
    sample: Optional[str]
    examples: List[str] = field(default_factory=list)
    metadata: Dict[str, Any] = field(default_factory=dict)
    origin: str = "row"


@dataclass
class ExcelSummary:
    sheets: List[str]
    metrics: Dict[str, float]
    column_samples: Dict[str, List[str]] = field(default_factory=dict)
    column_value_samples: Dict[str, Dict[str, List[str]]] = field(default_factory=dict)
    objects: List[ExcelObjectDescriptor] = field(default_factory=list)


@dataclass
class PresentationSummary:
    slide_count: int
    titles: List[str]


@dataclass
class WordField:
    label: str
    placeholder: str
    classification: str
    location: str
    endnote_ids: List[str] = field(default_factory=list)
    endnote_texts: List[str] = field(default_factory=list)


@dataclass
class PlaceholderMatch:
    text: str
    endnote_ids: List[str] = field(default_factory=list)


@dataclass
class WordSummary:
    headings: List[str]
    paragraphs: List[str]
    fields: List[WordField] = field(default_factory=list)
    stats: Dict[str, Any] = field(default_factory=dict)


@dataclass
class SwaggerOperationDescriptor:
    """Single API operation with its JSON attributes."""
    api_name: str
    endpoint: str
    method: str
    operation_id: Optional[str]
    summary: Optional[str]
    description: Optional[str]
    input_attributes: List[ExcelObjectDescriptor] = field(default_factory=list)
    output_attributes: List[ExcelObjectDescriptor] = field(default_factory=list)
    tags: List[str] = field(default_factory=list)


@dataclass
class SwaggerSummary:
    """Aggregated Swagger/OpenAPI specification."""
    api_title: str
    api_version: str
    api_description: Optional[str]
    operations: List[SwaggerOperationDescriptor]
    total_endpoints: int
    total_operations: int
    metrics: Dict[str, float] = field(default_factory=dict)


def _ensure_prance() -> None:
    global ResolvingParser
    if ResolvingParser is not None:
        return
    try:
        import importlib
        prance_module = importlib.import_module("prance")
        ResolvingParser = prance_module.ResolvingParser  # type: ignore[assignment]
    except ImportError as exc:  # pragma: no cover
        raise MappingDataError("prance is not installed; cannot parse Swagger/OpenAPI specs.") from exc


def _ensure_pandas() -> None:
    global pd
    if pd is not None:
        return
    try:
        import importlib

        pd = importlib.import_module("pandas")  # type: ignore[assignment]
    except ImportError as exc:  # pragma: no cover
        raise MappingDataError("pandas is not installed; cannot parse Excel files.") from exc


def _parse_json_as_excel(path: str) -> ExcelSummary:
    """Parse JSON data dictionary file and convert to ExcelSummary format."""
    logger.info("[mapping.parsers] Parsing JSON | path=%s", path)
    
    with open(path, 'r', encoding='utf-8') as f:
        data = json.load(f)
    
    # Detect structure: array of fields vs nested object
    fields_list = []
    
    if isinstance(data, dict):
        # Check for common data dictionary keys
        if 'fields' in data and isinstance(data['fields'], list):
            fields_list = data['fields']
        elif 'columns' in data and isinstance(data['columns'], list):
            fields_list = data['columns']
        elif 'attributes' in data and isinstance(data['attributes'], list):
            fields_list = data['attributes']
        else:
            # Flatten top-level dict keys as fields
            for key, value in data.items():
                if isinstance(value, dict):
                    fields_list.append({'fieldName': key, **value})
                else:
                    fields_list.append({'fieldName': key, 'value': value})
    elif isinstance(data, list):
        fields_list = data
    
    # Group fields by category (if present) to create "sheets"
    category_fields: Dict[str, List[Dict[str, Any]]] = {}
    
    for field in fields_list:
        if not isinstance(field, dict):
            continue
        
        category = field.get('category') or field.get('group') or field.get('section') or 'General'
        if category not in category_fields:
            category_fields[category] = []
        category_fields[category].append(field)
    
    # Build ExcelSummary structure
    sheets = list(category_fields.keys())
    metrics: Dict[str, float] = {}
    column_samples: Dict[str, List[str]] = {}
    column_value_samples: Dict[str, Dict[str, List[str]]] = {}
    object_descriptors: List[ExcelObjectDescriptor] = []
    
    for category, fields in category_fields.items():
        # Extract field names as "columns"
        field_names = []
        for idx, field in enumerate(fields):
            field_name = (
                field.get('fieldName') or 
                field.get('name') or 
                field.get('field') or 
                field.get('attribute') or 
                field.get('column') or 
                f'Field{idx+1}'
            )
            field_names.append(str(field_name))
            
            # Create descriptor for each field
            description = (
                field.get('description') or 
                field.get('desc') or 
                field.get('details') or 
                field.get('definition') or 
                ''
            )
            
            data_type = field.get('dataType') or field.get('type') or field.get('fieldType') or 'string'
            sample_value = field.get('sampleValue') or field.get('example') or field.get('default')
            
            # Extract metadata (everything except standard fields)
            metadata = {k: v for k, v in field.items() 
                       if k not in ('fieldName', 'name', 'description', 'desc', 'category', 'group')}
            
            descriptor = ExcelObjectDescriptor(
                sheet=category,
                row_index=idx,
                column='fieldName',
                name=field_name,
                description=description,
                path=f'$.{field_name}',
                sample=str(sample_value) if sample_value else None,
                examples=[str(sample_value)] if sample_value else [],
                metadata=metadata,
                origin='json'
            )
            object_descriptors.append(descriptor)
        
        column_samples[category] = field_names[:25]  # Limit to 25 samples
        
        # Create value samples from descriptions
        value_samples: Dict[str, List[str]] = {}
        for field in fields[:10]:  # Sample first 10 fields
            field_name = (
                field.get('fieldName') or 
                field.get('name') or 
                'Field'
            )
            desc = field.get('description', '')
            if desc:
                value_samples[str(field_name)] = [str(desc)[:160]]
        column_value_samples[category] = value_samples
        
        # Metrics: count of fields per category
        metrics[category] = float(len(fields))
    
    logger.info(
        "[mapping.parsers] JSON parsed | categories=%s total_fields=%s",
        len(sheets),
        len(object_descriptors)
    )
    
    return ExcelSummary(
        sheets=sheets,
        metrics=metrics,
        column_samples=column_samples,
        column_value_samples=column_value_samples,
        objects=object_descriptors,
    )


def parse_excel(path: str) -> ExcelSummary:
    if not os.path.exists(path):
        raise MappingDataError(f"Excel file missing: {path}")
    
    # Handle JSON files separately
    if path.lower().endswith('.json'):
        return _parse_json_as_excel(path)
    
    _ensure_pandas()
    
    # Handle CSV files separately
    if path.lower().endswith('.csv'):
        logger.info("[mapping.parsers] Parsing CSV | path=%s", path)
        df = pd.read_csv(path)  # type: ignore[attr-defined]
        sheet_name = os.path.basename(path).replace('.csv', '')
        
        sheets = [sheet_name]
        metrics: Dict[str, float] = {}
        column_samples: Dict[str, List[str]] = {}
        column_value_samples: Dict[str, Dict[str, List[str]]] = {}
        object_descriptors: List[ExcelObjectDescriptor] = []
        
        numeric_cols = df.select_dtypes(include=["number"])  # type: ignore[arg-type]
        if not numeric_cols.empty:
            metrics[sheet_name] = float(numeric_cols.sum().sum())
        column_samples[sheet_name] = [str(col) for col in list(df.columns)[:25]]
        column_value_samples[sheet_name] = _sample_column_values(df)
        object_descriptors.extend(_extract_excel_objects(sheet_name, df))
        
        if object_descriptors:
            logger.info(
                "[mapping.parsers] CSV object descriptors extracted | rows=%s objects=%s",
                len(df),
                len(object_descriptors),
            )
        
        return ExcelSummary(
            sheets=sheets,
            metrics=metrics,
            column_samples=column_samples,
            column_value_samples=column_value_samples,
            objects=object_descriptors,
        )
    
    # Handle Excel files (.xlsx, .xls)
    logger.info("[mapping.parsers] Parsing Excel | path=%s", path)
    workbook = pd.ExcelFile(path)  # type: ignore[attr-defined]
    sheets = [str(name) for name in workbook.sheet_names]
    metrics: Dict[str, float] = {}
    column_samples: Dict[str, List[str]] = {}
    column_value_samples: Dict[str, Dict[str, List[str]]] = {}
    object_descriptors: List[ExcelObjectDescriptor] = []
    for name in sheets:
        df = workbook.parse(name)
        numeric_cols = df.select_dtypes(include=["number"])  # type: ignore[arg-type]
        if not numeric_cols.empty:
            metrics[str(name)] = float(numeric_cols.sum().sum())
        column_samples[str(name)] = [str(col) for col in list(df.columns)[:25]]
        column_value_samples[str(name)] = _sample_column_values(df)
        object_descriptors.extend(_extract_excel_objects(str(name), df))
    if object_descriptors:
        logger.info(
            "[mapping.parsers] Excel object descriptors extracted | sheets=%s objects=%s",
            len(sheets),
            len(object_descriptors),
        )
    return ExcelSummary(
        sheets=sheets,
        metrics=metrics,
        column_samples=column_samples,
        column_value_samples=column_value_samples,
        objects=object_descriptors,
    )


def _sample_column_values(df: Any) -> Dict[str, List[str]]:
    samples: Dict[str, List[str]] = {}
    for column in df.columns:
        values: List[str] = []
        for value in df[column]:
            text = _clean_text(value)
            if text:
                values.append(text[:160])
            if len(values) >= 3:
                break
        if values:
            samples[str(column)] = values
    return samples


def _safe_row_index(index_value: Any, fallback: int) -> int:
    try:
        if isinstance(index_value, (int, float)):
            if pd is not None and pd.isna(index_value):  # type: ignore[attr-defined]
                return fallback
            return int(index_value)
    except Exception:  # pragma: no cover
        return fallback
    return fallback


def _extract_excel_objects(sheet_name: str, df: Any) -> List[ExcelObjectDescriptor]:
    descriptors: List[ExcelObjectDescriptor] = []
    seen: set = set()
    max_objects = _DESCRIPTOR_LIMIT_PER_SHEET
    for row_pos, (idx, row) in enumerate(df.iterrows()):
        if len(descriptors) >= max_objects:
            break
        row_index = _safe_row_index(idx, row_pos)
        row_mapping = {str(col): row[col] for col in df.columns}
        row_descriptor = _descriptor_from_row(row_mapping)
        if row_descriptor:
            _register_descriptor(
                descriptors,
                seen,
                sheet_name,
                row_index,
                None,
                row_descriptor,
                origin="row",
            )
        for column in df.columns:
            cell_value = row[column]
            payloads = _extract_json_payloads(cell_value)
            for payload in payloads:
                descriptor_payload = _descriptor_from_mapping(payload)
                if descriptor_payload:
                    _register_descriptor(
                        descriptors,
                        seen,
                        sheet_name,
                        row_index,
                        str(column),
                        descriptor_payload,
                        origin="json",
                    )
                if len(descriptors) >= max_objects:
                    break
            if len(descriptors) >= max_objects:
                break
    return descriptors


def _descriptor_from_row(row_mapping: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    normalized = {str(key): row_mapping[key] for key in row_mapping}
    name = _extract_value(normalized, _FIELD_NAME_HINTS)
    description = _extract_value(normalized, _FIELD_DESCRIPTION_HINTS)
    path = _extract_value(normalized, _FIELD_PATH_HINTS)
    sample = _extract_value(normalized, _FIELD_SAMPLE_HINTS)
    
    # Must have either a name or path
    if not (name or path):
        return None
    
    # For ACORD/XML path CSVs: if we have a path but no description, use the path as description
    if path and (not description or len(description) < 4):
        description = f"Path: {path}"
    
    # Still require some description after the above fallback
    if not description or len(description) < 4:
        return None
    
    metadata = _filter_metadata(normalized)
    return {
        "name": name or path,
        "description": description,
        "path": path or name,
        "sample": sample,
        "examples": [],
        "metadata": metadata,
    }


def _extract_json_payloads(value: Any) -> List[Dict[str, Any]]:
    text = _clean_text(value)
    if not text or text[0] not in "{[":
        return []
    try:
        parsed = json.loads(text)
    except json.JSONDecodeError:
        return []
    if isinstance(parsed, dict):
        return [parsed]
    if isinstance(parsed, list):
        return [item for item in parsed if isinstance(item, dict)]
    return []


def _descriptor_from_mapping(payload: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    if not isinstance(payload, dict):
        return None
    name = _extract_value(payload, _FIELD_NAME_HINTS)
    description = _extract_value(payload, _FIELD_DESCRIPTION_HINTS)
    path = _extract_value(payload, _FIELD_PATH_HINTS)
    sample = _extract_value(payload, _FIELD_SAMPLE_HINTS)
    raw_examples = payload.get("examples") or payload.get("samples")
    examples: List[str] = []
    if isinstance(raw_examples, list):
        for entry in raw_examples:
            text = _clean_text(entry)
            if text:
                examples.append(text)
    if not description or len(description) < 4:
        return None
    if not (name or path):
        return None
    metadata = _filter_metadata(payload)
    return {
        "name": name or path,
        "description": description,
        "path": path or name,
        "sample": sample,
        "examples": examples,
        "metadata": metadata,
    }


def _register_descriptor(
    descriptors: List[ExcelObjectDescriptor],
    seen: set,
    sheet: str,
    row_index: int,
    column: Optional[str],
    payload: Dict[str, Any],
    origin: str,
) -> None:
    key = (
        sheet,
        payload.get("path") or payload.get("name"),
        payload.get("description"),
    )
    if key in seen:
        return
    descriptor = ExcelObjectDescriptor(
        sheet=sheet,
        row_index=row_index,
        column=column,
        name=payload.get("name"),
        description=payload.get("description"),
        path=payload.get("path"),
        sample=payload.get("sample"),
        examples=list(payload.get("examples", [])),
        metadata=payload.get("metadata", {}),
        origin=origin,
    )
    descriptors.append(descriptor)
    seen.add(key)


def _unique_in_order(values: Sequence[str]) -> List[str]:
    seen: set[str] = set()
    ordered: List[str] = []
    for value in values:
        if value in seen:
            continue
        seen.add(value)
        ordered.append(value)
    return ordered


def _extract_endnote_refs_from_run(run: Any) -> List[str]:
    refs: List[str] = []
    element = getattr(run, "_element", None)
    if element is None:
        return refs
    for node in element.findall(".//" + _WORD_NAMESPACE + "endnoteReference"):
        ref_id = node.get(_WORD_NAMESPACE + "id")
        if ref_id:
            refs.append(ref_id)
    return refs


def _collect_following_endnote_ids(runs: Sequence[Any], run_index: int) -> List[str]:
    ids: List[str] = []
    for pointer in range(run_index + 1, len(runs)):
        run = runs[pointer]
        refs = _extract_endnote_refs_from_run(run)
        if refs:
            ids.extend(refs)
        text = getattr(run, "text", "") or ""
        if text:
            break
        if refs:
            continue
    return ids


def _extract_bracket_placeholders_from_runs(runs: Optional[Sequence[Any]]) -> List[PlaceholderMatch]:
    matches: List[PlaceholderMatch] = []
    if not runs:
        return matches
    buffer: List[str] = []
    collecting = False
    for index, run in enumerate(runs):
        text = getattr(run, "text", "") or ""
        if not text:
            continue
        for char in text:
            if char == "[":
                if collecting and buffer:
                    token = "[" + "".join(buffer).strip()
                    if len(token) > 1:
                        matches.append(PlaceholderMatch(text=token + "]"))
                collecting = True
                buffer = []
                continue
            if char == "]" and collecting:
                token = "[" + "".join(buffer).strip() + "]"
                if len(token) > 2:
                    endnote_ids = _collect_following_endnote_ids(runs, index)
                    matches.append(PlaceholderMatch(text=token, endnote_ids=endnote_ids))
                buffer = []
                collecting = False
                continue
            if collecting:
                buffer.append(char)
        if "\n" in text or "\r" in text:
            if collecting and buffer:
                token = "[" + "".join(buffer).strip()
                if len(token) > 1:
                    matches.append(PlaceholderMatch(text=token))
                buffer = []
                collecting = False
    if collecting and buffer:
        token = "[" + "".join(buffer).strip()
        if len(token) > 1:
            matches.append(PlaceholderMatch(text=token))
    return matches


def _extract_placeholders(text: str, runs: Optional[Sequence[Any]] = None) -> List[PlaceholderMatch]:
    matches: List[PlaceholderMatch] = []
    if text:
        for token in _ANGLE_PLACEHOLDER_PATTERN.findall(text):
            normalized = token.strip()
            if normalized:
                matches.append(PlaceholderMatch(text=normalized))
    run_matches = _extract_bracket_placeholders_from_runs(runs)
    if run_matches:
        matches.extend(run_matches)
    elif text:
        for token in _BRACKET_PLACEHOLDER_PATTERN.findall(text):
            normalized = token.strip()
            if normalized and normalized != "[]":
                matches.append(PlaceholderMatch(text=normalized))
        if "[" in text and "]" not in text:
            stripped = text.strip()
            if stripped.startswith("[") and len(stripped) > 1:
                matches.append(PlaceholderMatch(text=stripped))
    return matches


def _collect_cell_runs(cell: Any) -> List[Any]:
    runs: List[Any] = []
    for paragraph in getattr(cell, "paragraphs", []) or []:
        runs.extend(getattr(paragraph, "runs", []) or [])
    return runs


def _load_endnote_definitions(path: str) -> Dict[str, str]:
    try:
        with zipfile.ZipFile(path) as archive:
            if _ENDNOTE_PATH not in archive.namelist():
                return {}
            data = archive.read(_ENDNOTE_PATH)
    except (FileNotFoundError, zipfile.BadZipFile):
        logger.warning("[mapping.parsers] Unable to read endnotes | path=%s", path)
        return {}
    except Exception as exc:  # pragma: no cover
        logger.warning("[mapping.parsers] Unexpected endnote error | path=%s error=%s", path, exc)
        return {}
    definitions: Dict[str, str] = {}
    try:
        root = ET.fromstring(data)
    except ET.ParseError:  # pragma: no cover
        logger.warning("[mapping.parsers] Invalid endnotes xml | path=%s", path)
        return {}
    for endnote in root.findall(f".//{_WORD_NAMESPACE}endnote"):
        note_id = endnote.get(f"{_WORD_NAMESPACE}id")
        if not note_id or note_id in {"-1", "0"}:
            continue
        paragraphs: List[str] = []
        for paragraph in endnote.findall(f"{_WORD_NAMESPACE}p"):
            texts = [node.text or "" for node in paragraph.findall(f".//{_WORD_NAMESPACE}t")]
            paragraph_text = "".join(texts).strip()
            if paragraph_text:
                paragraphs.append(paragraph_text)
        if paragraphs:
            definitions[note_id] = "\n".join(paragraphs)
    return definitions


def _extract_superscript_digits(text: str) -> List[str]:
    digits: List[str] = []
    for char in text or "":
        digit = _SUPERSCRIPT_MAP.get(char)
        if digit is not None:
            digits.append(digit)
    return digits


def _extract_run_superscripts(runs: Optional[Sequence[Any]]) -> List[str]:
    if not runs:
        return []
    digits: List[str] = []
    buffer = ""
    collecting = False
    for run in runs:
        text = getattr(run, "text", "") or ""
        if not text:
            continue
        is_super = bool(getattr(getattr(run, "font", None), "superscript", False))
        if is_super:
            collecting = True
            buffer += text
            continue
        if collecting and buffer:
            digits.extend(_SUPERSCRIPT_DIGIT_PATTERN.findall(buffer))
            buffer = ""
            collecting = False
    if collecting and buffer:
        digits.extend(_SUPERSCRIPT_DIGIT_PATTERN.findall(buffer))
    return digits


def _extract_cell_superscripts(cell: Any) -> List[str]:
    digits: List[str] = []
    for paragraph in getattr(cell, "paragraphs", []) or []:
        digits.extend(_extract_run_superscripts(getattr(paragraph, "runs", None)))
    return digits


def _extract_superscript_definitions(paragraphs: Sequence[Any]) -> tuple[Dict[str, str], set[int]]:
    definitions: Dict[str, str] = {}
    indexes: set[int] = set()
    for idx, paragraph in enumerate(paragraphs):
        raw_text = getattr(paragraph, "text", "") or ""
        stripped = raw_text.strip()
        if not stripped:
            continue
        match = _SUPERSCRIPT_DEF_PATTERN.match(stripped)
        if match:
            digit = _SUPERSCRIPT_MAP.get(match.group("sup"))
            if digit is None:
                continue
            body = match.group("body").strip(" :-\u2013\u2014")
            if not body:
                continue
            definitions.setdefault(digit, body)
            indexes.add(idx)
            continue
        digit = _extract_superscript_prefix_digit(paragraph)
        if not digit or digit in definitions:
            continue
        body = _strip_definition_prefix(stripped, digit)
        if not body:
            continue
        definitions[digit] = body
        indexes.add(idx)
    return definitions, indexes


def _strip_definition_prefix(text: str, digit: str) -> str:
    pattern = re.compile(rf"^\s*[\[\(]*{re.escape(digit)}(?:st|nd|rd|th)?[\]\)\.\-:\u2013\u2014\s]*")
    return pattern.sub("", text, count=1).strip()


def _extract_superscript_prefix_digit(paragraph: Any) -> Optional[str]:
    runs = getattr(paragraph, "runs", None)
    if not runs:
        return None
    buffer: List[str] = []
    seen_super = False
    for run in runs:
        text = getattr(run, "text", "") or ""
        if not text and not seen_super:
            continue
        is_super = bool(getattr(getattr(run, "font", None), "superscript", False))
        if not is_super:
            break
        seen_super = True
        buffer.append(text)
    if not seen_super:
        return None
    candidate = "".join(buffer).strip()
    if not candidate:
        return None
    match = _SUPERSCRIPT_DIGIT_PATTERN.search(candidate)
    return match.group(0) if match else None


def _format_superscript_placeholder(label: Optional[str], digit: str) -> str:
    token = (label or f"superscript_{digit}").strip()
    token = re.sub(r"\s+", " ", token)
    core = token.strip("<>") or f"superscript_{digit}"
    return f"<{core}>"


def parse_presentation(path: str) -> PresentationSummary:
    if not os.path.exists(path):
        raise MappingDataError(f"Presentation file missing: {path}")
    _ensure_pptx()
    logger.info("[mapping.parsers] Parsing presentation | path=%s", path)
    presentation = pptx.Presentation(path)  # type: ignore[attr-defined]
    titles: List[str] = []
    for slide in presentation.slides:
        if slide.shapes.title and slide.shapes.title.text:
            titles.append(slide.shapes.title.text.strip())
    return PresentationSummary(slide_count=len(presentation.slides), titles=titles)


def _classify_paragraph(style_name: str, text: str) -> str:
    if "heading" in style_name:
        return "section_heading"
    if text.endswith(":"):
        return "label"
    return "paragraph"


def parse_word_document(path: str) -> WordSummary:
    if not os.path.exists(path):
        raise MappingDataError(f"Word file missing: {path}")
    _ensure_docx()
    logger.info("[mapping.parsers] Parsing word doc | path=%s", path)
    document = docx.Document(path)  # type: ignore[attr-defined]
    doc_paragraphs = list(document.paragraphs)
    superscript_definitions, superscript_definition_indexes = _extract_superscript_definitions(doc_paragraphs)
    superscript_usage: Counter[str] = Counter()
    endnote_definitions = _load_endnote_definitions(path)
    endnote_usage: Counter[str] = Counter()
    headings: List[str] = []
    paragraphs: List[str] = []
    fields: List[WordField] = []
    static_sections: List[str] = []
    dynamic_count = 0
    location_index = 0
    current_section = "Document Start"  # Track current section heading

    def _record_field(
        label: str,
        placeholder: str,
        classification: str,
        location_hint: str,
        endnote_ids: Optional[List[str]] = None,
        section_heading: Optional[str] = None,
    ) -> None:
        nonlocal dynamic_count
        note_ids = [str(identifier) for identifier in (endnote_ids or []) if identifier]
        note_texts = [endnote_definitions.get(identifier, "") for identifier in note_ids if endnote_definitions.get(identifier)]
        dynamic_count += 1
        for identifier in note_ids:
            endnote_usage[identifier] += 1
        
        # Enhanced location with section heading
        enhanced_location = location_hint
        if section_heading:
            enhanced_location = f"Section: {section_heading} | {location_hint}"
            
        fields.append(
            WordField(
                label=label,
                placeholder=placeholder,
                classification=classification,
                location=enhanced_location,
                endnote_ids=note_ids,
                endnote_texts=note_texts,
            )
        )

    def _record_superscript_field(label: str, digit: str, classification: str, location_hint: str, section_heading: Optional[str] = None) -> None:
        placeholder_token = _format_superscript_placeholder(superscript_definitions.get(digit), digit)
        _record_field(label, placeholder_token, classification, location_hint, section_heading=section_heading)
        superscript_usage[digit] += 1

    for p_index, paragraph in enumerate(doc_paragraphs):
        text = paragraph.text.strip()
        if not text:
            continue
        style_name = getattr(paragraph.style, "name", "").lower()
        paragraph_runs = getattr(paragraph, "runs", None)
        placeholders = _extract_placeholders(text, paragraph_runs)
        superscript_digits = _extract_superscript_digits(text)
        superscript_digits.extend(_extract_run_superscripts(paragraph_runs))
        classification = _classify_paragraph(style_name, text)
        if "heading" in style_name:
            headings.append(text)
            current_section = text  # Update current section
            if not placeholders and not superscript_digits:
                static_sections.append(text)
        paragraphs.append(text)
        if p_index in superscript_definition_indexes:
            continue
        if placeholders:
            for match in placeholders:
                location_index += 1
                _record_field(text, match.text, classification, f"paragraph[{location_index}]", match.endnote_ids, section_heading=current_section)
        if superscript_digits:
            for digit in _unique_in_order(superscript_digits):
                location_index += 1
                _record_superscript_field(text, digit, classification, f"paragraph[{location_index}]", section_heading=current_section)

    for t_index, table in enumerate(document.tables):  # type: ignore[attr-defined]
        for r_index, row in enumerate(table.rows):
            for c_index, cell in enumerate(row.cells):
                cell_text = cell.text.strip()
                if not cell_text:
                    continue
                cell_runs = _collect_cell_runs(cell)
                placeholders = _extract_placeholders(cell_text, cell_runs)
                superscript_digits = _extract_superscript_digits(cell_text)
                superscript_digits.extend(_extract_cell_superscripts(cell))
                if not placeholders and not superscript_digits:
                    continue
                paragraphs.append(cell_text)
                for match in placeholders:
                    location_index += 1
                    _record_field(
                        cell_text,
                        match.text,
                        "table_cell",
                        f"table[{t_index+1}] row[{r_index+1}] col[{c_index+1}]",
                        match.endnote_ids,
                        section_heading=current_section,
                    )
                if superscript_digits:
                    for digit in _unique_in_order(superscript_digits):
                        location_index += 1
                        _record_superscript_field(
                            cell_text,
                            digit,
                            "table_cell",
                            f"table[{t_index+1}] row[{r_index+1}] col[{c_index+1}]",
                            section_heading=current_section,
                        )

    stats = {
        "dynamic_fields": dynamic_count,
        "static_sections": static_sections,
        "total_paragraphs": len(paragraphs),
        "superscript_definitions": superscript_definitions,
        "superscript_definition_count": len(superscript_definitions),
        "superscript_fields_detected": sum(superscript_usage.values()),
        "superscript_usage": dict(superscript_usage),
        "endnote_definitions": endnote_definitions,
        "endnote_definition_count": len(endnote_definitions),
        "endnote_fields_detected": sum(endnote_usage.values()),
        "endnote_usage": dict(endnote_usage),
    }
    unmatched_definitions = sorted(set(superscript_definitions.keys()) - set(superscript_usage.keys()))
    if unmatched_definitions:
        stats["superscript_unmatched_definitions"] = unmatched_definitions
    unmatched_endnotes = sorted(set(endnote_definitions.keys()) - set(endnote_usage.keys()))
    if unmatched_endnotes:
        stats["endnote_unmatched_definitions"] = unmatched_endnotes
    logger.info(
        "[mapping.parsers] Word summary | dynamic_fields=%s static_sections=%s superscript_fields=%s",
        dynamic_count,
        len(static_sections),
        stats.get("superscript_fields_detected", 0),
    )
    return WordSummary(headings=headings, paragraphs=paragraphs, fields=fields, stats=stats)


def _repeat_distribution(counter: Counter) -> List[Dict[str, int]]:
    repeat_counter = Counter(counter.values())
    ordered: List[Dict[str, int]] = []
    for occurrences, placeholder_count in sorted(repeat_counter.items(), key=lambda item: item[0], reverse=True):
        ordered.append({"occurrences": occurrences, "placeholder_count": placeholder_count})
    return ordered


def _classification_breakdown(fields: List[WordField]) -> Dict[str, int]:
    breakdown = Counter((field.classification or "unknown") for field in fields)
    return dict(breakdown)


def _detect_superscripts(paragraphs: List[str]) -> Dict[str, int]:
    counts: Counter = Counter()
    for text in paragraphs or []:
        for char in text:
            if char in _SUPERSCRIPT_CHARS:
                digit = _SUPERSCRIPT_MAP.get(char)
                if digit is not None:
                    counts[digit] += 1
    return dict(counts)


def summarize_word_analysis(summary: WordSummary) -> Dict[str, Any]:
    placeholder_counts = Counter((field.placeholder or "").strip() for field in summary.fields if field.placeholder)
    placeholder_counts.pop("", None)
    classification_counts = _classification_breakdown(summary.fields)
    table_fields = classification_counts.get("table_cell", 0)
    non_table_fields = sum(classification_counts.values()) - table_fields
    repeat_distribution = _repeat_distribution(placeholder_counts)
    raw_superscripts = _detect_superscripts(summary.paragraphs)
    superscript_definitions: Dict[str, str] = summary.stats.get("superscript_definitions", {}) or {}
    superscript_usage: Dict[str, int] = summary.stats.get("superscript_usage", {}) or {}
    unmatched_definitions: List[str] = summary.stats.get("superscript_unmatched_definitions", []) or []
    definition_count = summary.stats.get("superscript_definition_count", len(superscript_definitions))
    superscript_fields_detected = summary.stats.get("superscript_fields_detected", 0)
    endnote_definitions: Dict[str, str] = summary.stats.get("endnote_definitions", {}) or {}
    endnote_usage: Dict[str, int] = summary.stats.get("endnote_usage", {}) or {}
    endnote_unmatched: List[str] = summary.stats.get("endnote_unmatched_definitions", []) or []
    endnote_definition_count = summary.stats.get("endnote_definition_count", len(endnote_definitions))
    endnote_fields_detected = summary.stats.get("endnote_fields_detected", 0)
    unresolved_digits = sorted(set(raw_superscripts.keys()) - set(superscript_definitions.keys()))
    total_placeholders = sum(placeholder_counts.values())
    unique_placeholders = len(placeholder_counts)
    insights: List[str] = []
    if total_placeholders:
        insights.append(
            f"Identified {unique_placeholders} unique placeholders across {total_placeholders} total occurrences."
        )
    if repeat_distribution:
        top = repeat_distribution[0]
        if top["occurrences"] > 1:
            insights.append(
                f"{top['placeholder_count']} placeholders repeat {top['occurrences']} times, suggesting multi-instance blocks."
            )
    if superscript_fields_detected:
        insights.append(
            f"Synthesized {superscript_fields_detected} superscript-driven placeholders across {definition_count} definitions."
        )
    elif superscript_definitions:
        insights.append(
            f"Detected {definition_count} superscript definition entries but no references inside the main body."
        )
    elif raw_superscripts:
        portions = ", ".join(f"{digit} (x{count})" for digit, count in sorted(raw_superscripts.items()))
        insights.append(f"Superscript annotations detected for: {portions}.")
    if endnote_fields_detected:
        insights.append(
            f"Linked {endnote_fields_detected} placeholders to {endnote_definition_count} disclosure endnotes."
        )
    elif endnote_definitions:
        insights.append(
            f"Detected {endnote_definition_count} endnote definitions but no linked placeholders."
        )
    insights.append(
        f"Paragraph placeholders: {non_table_fields}, table placeholders: {table_fields}, static sections: {len(summary.stats.get('static_sections', []))}."
    )
    superscript_details: List[str] = []
    for digit, body in sorted(superscript_definitions.items(), key=lambda item: item[0]):
        usage = superscript_usage.get(digit, 0)
        status = "referenced" if usage else "not referenced"
        superscript_details.append(f"{digit} → {body.strip()} ({status})")
    if unmatched_definitions:
        superscript_details.append(
            f"Definitions without in-body references: {', '.join(unmatched_definitions)}"
        )
    if unresolved_digits:
        superscript_details.append(
            f"Superscripts missing definition entries: {', '.join(unresolved_digits)}"
        )
    superscripts = {
        "total": definition_count,
        "missing": len(unmatched_definitions),
        "unresolved": len(unresolved_digits),
        "details": superscript_details,
    }
    endnote_details: List[str] = []
    def _endnote_sort_key(item: tuple[str, str]) -> tuple[int, str]:
        identifier, text = item
        if identifier.isdigit():
            return (int(identifier), identifier)
        return (10**6, identifier)

    for note_id, body in sorted(endnote_definitions.items(), key=_endnote_sort_key):
        usage = endnote_usage.get(note_id, 0)
        status = "referenced" if usage else "not referenced"
        endnote_details.append(f"{note_id} → {body.strip()} ({status})")
    if endnote_unmatched:
        endnote_details.append(
            f"Endnotes without field references: {', '.join(endnote_unmatched)}"
        )
    endnotes = {
        "total": endnote_definition_count,
        "missing": len(endnote_unmatched),
        "details": endnote_details,
    }
    return {
        "totals": {
            "paragraphs": len(summary.paragraphs),
            "placeholders": total_placeholders,
            "unique_placeholders": unique_placeholders,
            "table_fields": table_fields,
            "non_table_fields": non_table_fields,
            "static_sections": len(summary.stats.get("static_sections", [])),
        },
        "repeatDistribution": repeat_distribution,
        "classificationBreakdown": classification_counts,
        "superscripts": superscripts,
        "endnotes": endnotes,
        "insights": insights,
    }


def _descriptor_is_array(descriptor: ExcelObjectDescriptor) -> bool:
    path = (descriptor.path or "").lower()
    name = (descriptor.name or "").lower()
    metadata_values = " ".join(str(value).lower() for value in descriptor.metadata.values()) if descriptor.metadata else ""
    if "[]" in path or "[]" in name:
        return True
    if any(token in path for token in ("[0]", "[1]", ".items")):
        return True
    if "array" in metadata_values or "list" in metadata_values:
        return True
    if descriptor.examples and len(descriptor.examples) > 1:
        return True
    return False


def summarize_excel_analysis(summary: ExcelSummary) -> Dict[str, Any]:
    column_counter: Counter = Counter()
    sheet_overview: List[Dict[str, Any]] = []
    for sheet, columns in summary.column_samples.items():
        column_counter.update(column.strip() for column in columns if column)
        sheet_overview.append(
            {
                "sheet": sheet,
                "column_count": len(columns),
                "sampled_columns": len(summary.column_value_samples.get(sheet, {})),
                "object_descriptors": len([obj for obj in summary.objects if obj.sheet == sheet]),
            }
        )
    repeat_distribution = _repeat_distribution(column_counter)
    descriptor_total = len(summary.objects)
    array_descriptors = [descriptor for descriptor in summary.objects if _descriptor_is_array(descriptor)]
    arrays_by_sheet = Counter(descriptor.sheet for descriptor in array_descriptors)
    total_columns = sum(column_counter.values())
    unique_columns = len(column_counter)
    insights: List[str] = []
    if summary.sheets:
        insights.append(
            f"Parsed {len(summary.sheets)} sheets with {unique_columns} unique column labels ({total_columns} total occurrences)."
        )
    if repeat_distribution:
        top = repeat_distribution[0]
        if top["occurrences"] > 1:
            insights.append(
                f"{top['placeholder_count']} column labels repeat {top['occurrences']} times across sheets."
            )
    if descriptor_total:
        insights.append(
            f"Extracted {descriptor_total} JSON/object descriptors; {len(array_descriptors)} appear to be arrays."
        )
    if arrays_by_sheet:
        sheet_details = ", ".join(f"{sheet}: {count}" for sheet, count in arrays_by_sheet.items())
        insights.append("Array descriptors by sheet → " + sheet_details)
    return {
        "totals": {
            "sheets": len(summary.sheets),
            "total_columns": total_columns,
            "unique_columns": unique_columns,
            "json_descriptors": descriptor_total,
            "array_descriptors": len(array_descriptors),
        },
        "repeatDistribution": repeat_distribution,
        "sheetOverview": sheet_overview,
        "arraysBySheet": dict(arrays_by_sheet),
        "insights": insights,
    }


def parse_swagger(path: str) -> SwaggerSummary:
    """Parse Swagger/OpenAPI specification file with automatic $ref resolution.
    
    Args:
        path: Path to Swagger/OpenAPI file (.json, .yaml, .yml)
        
    Returns:
        SwaggerSummary with all operations and flattened JSON schemas
        
    Raises:
        MappingDataError: If file doesn't exist, parsing fails, or prance not installed
    """
    if not os.path.exists(path):
        raise MappingDataError(f"Swagger file missing: {path}")
    
    _ensure_prance()
    
    logger.info("[mapping.parsers] Parsing Swagger/OpenAPI | path=%s", path)
    
    try:
        # Parse with automatic $ref resolution
        parser = ResolvingParser(path)  # type: ignore[misc]
        spec = parser.specification
    except Exception as exc:
        logger.error("[mapping.parsers] Failed to parse Swagger spec | error=%s", exc, exc_info=True)
        raise MappingDataError(f"Failed to parse Swagger/OpenAPI spec: {exc}") from exc
    
    # Extract metadata
    api_title = spec.get('info', {}).get('title', 'Unnamed API')
    api_version = spec.get('info', {}).get('version', '1.0')
    api_description = spec.get('info', {}).get('description')
    
    # Parse all operations
    operations: List[SwaggerOperationDescriptor] = []
    paths = spec.get('paths', {})
    
    for endpoint, path_item in paths.items():
        for method in ['get', 'post', 'put', 'delete', 'patch', 'head', 'options']:
            if method not in path_item:
                continue
            
            operation = path_item[method]
            
            # Extract operation metadata
            operation_id = operation.get('operationId', f"{method}_{endpoint.replace('/', '_')}")
            summary = operation.get('summary', '')
            description = operation.get('description', '')
            tags = operation.get('tags', [])
            
            # Extract input schema (request body)
            input_attrs = _extract_swagger_request_attributes(operation, endpoint, method)
            
            # Extract output schema (response body)
            output_attrs = _extract_swagger_response_attributes(operation, endpoint, method)
            
            operations.append(SwaggerOperationDescriptor(
                api_name=api_title,
                endpoint=endpoint,
                method=method.upper(),
                operation_id=operation_id,
                summary=summary,
                description=description,
                input_attributes=input_attrs,
                output_attributes=output_attrs,
                tags=tags
            ))
    
    # Calculate metrics
    metrics = {
        'total_endpoints': len(paths),
        'total_operations': len(operations),
        'total_input_attributes': sum(len(op.input_attributes) for op in operations),
        'total_output_attributes': sum(len(op.output_attributes) for op in operations),
    }
    
    logger.info(
        "[mapping.parsers] Parsed Swagger spec | api=%s version=%s endpoints=%s operations=%s input_attrs=%s output_attrs=%s",
        api_title,
        api_version,
        len(paths),
        len(operations),
        metrics['total_input_attributes'],
        metrics['total_output_attributes']
    )
    
    # Log each operation's details
    for op in operations:
        logger.info(
            "[mapping.parsers] Operation detail | id=%s method=%s endpoint=%s inputs=%d outputs=%d",
            op.operation_id,
            op.method,
            op.endpoint,
            len(op.input_attributes),
            len(op.output_attributes)
        )
    
    return SwaggerSummary(
        api_title=api_title,
        api_version=api_version,
        api_description=api_description,
        operations=operations,
        total_endpoints=len(paths),
        total_operations=len(operations),
        metrics=metrics
    )


def _extract_swagger_request_attributes(
    operation: Dict[str, Any],
    endpoint: str,
    method: str
) -> List[ExcelObjectDescriptor]:
    """Extract request body JSON schema as flattened attributes."""
    
    # OpenAPI 3.x
    if 'requestBody' in operation:
        content = operation['requestBody'].get('content', {})
        json_content = content.get('application/json', {})
        schema = json_content.get('schema')
        if schema:
            return _flatten_json_schema(schema, f"{endpoint}::{method}::request", "request")
    
    # Swagger 2.0 - parameters with in=body
    for param in operation.get('parameters', []):
        if param.get('in') == 'body':
            schema = param.get('schema')
            if schema:
                return _flatten_json_schema(schema, f"{endpoint}::{method}::request", "request")
    
    return []


def _extract_swagger_response_attributes(
    operation: Dict[str, Any],
    endpoint: str,
    method: str
) -> List[ExcelObjectDescriptor]:
    """Extract 200 response JSON schema as flattened attributes."""
    responses = operation.get('responses', {})
    
    # Try 200, then 201, then 2xx
    for status_code in ['200', '201', '2XX', '2xx']:
        if status_code not in responses:
            continue
        
        response = responses[status_code]
        
        # OpenAPI 3.x
        content = response.get('content', {})
        json_content = content.get('application/json', {})
        if 'schema' in json_content:
            return _flatten_json_schema(
                json_content['schema'],
                f"{endpoint}::{method}::response",
                "response"
            )
        
        # Swagger 2.0
        if 'schema' in response:
            return _flatten_json_schema(
                response['schema'],
                f"{endpoint}::{method}::response",
                "response"
            )
    
    return []


def _flatten_json_schema(
    schema: Dict[str, Any],
    sheet: str,
    prefix: str = '',
    path: str = ''
) -> List[ExcelObjectDescriptor]:
    """Recursively flatten JSON schema to ExcelObjectDescriptor list.
    
    Args:
        schema: JSON schema dict (already resolved, no $ref)
        sheet: Sheet identifier (e.g., "endpoint::method::request")
        prefix: Prefix for origin tracking
        path: Current path in dot notation
        
    Returns:
        List of ExcelObjectDescriptor for all leaf and container attributes
    """
    descriptors: List[ExcelObjectDescriptor] = []
    schema_type = schema.get('type', 'object')
    
    # Handle object with properties
    if schema_type == 'object' or 'properties' in schema:
        properties = schema.get('properties', {})
        required_fields = schema.get('required', [])
        
        for prop_name, prop_schema in properties.items():
            current_path = f"{path}.{prop_name}" if path else prop_name
            prop_type = prop_schema.get('type', 'unknown')
            
            # Leaf element (primitive types)
            if prop_type in ['string', 'number', 'integer', 'boolean']:
                descriptors.append(ExcelObjectDescriptor(
                    sheet=sheet,
                    row_index=len(descriptors),
                    column=prop_name,
                    name=prop_name,
                    description=prop_schema.get('description', ''),
                    path=current_path,
                    sample=str(prop_schema.get('example', '')),
                    examples=[str(prop_schema.get('example', ''))] if 'example' in prop_schema else [],
                    metadata={
                        'type': prop_type,
                        'format': prop_schema.get('format', ''),
                        'enum': prop_schema.get('enum', []),
                        'required': prop_name in required_fields,
                        'is_leaf': True,
                        'is_array': False
                    },
                    origin=prefix
                ))
            
            # Array element
            elif prop_type == 'array':
                items_schema = prop_schema.get('items', {})
                
                # Add array container
                descriptors.append(ExcelObjectDescriptor(
                    sheet=sheet,
                    row_index=len(descriptors),
                    column=prop_name,
                    name=prop_name,
                    description=prop_schema.get('description', ''),
                    path=current_path,
                    sample='',
                    metadata={
                        'type': 'array',
                        'required': prop_name in required_fields,
                        'is_leaf': False,
                        'is_array': True
                    },
                    origin=prefix
                ))
                
                # Recurse into array items
                nested = _flatten_json_schema(
                    items_schema,
                    sheet,
                    prefix,
                    f"{current_path}[]"
                )
                descriptors.extend(nested)
            
            # Nested object
            elif prop_type == 'object' or 'properties' in prop_schema:
                # Add object container
                descriptors.append(ExcelObjectDescriptor(
                    sheet=sheet,
                    row_index=len(descriptors),
                    column=prop_name,
                    name=prop_name,
                    description=prop_schema.get('description', ''),
                    path=current_path,
                    sample='',
                    metadata={
                        'type': 'object',
                        'required': prop_name in required_fields,
                        'is_leaf': False,
                        'is_array': False
                    },
                    origin=prefix
                ))
                
                # Recurse into nested properties
                nested = _flatten_json_schema(
                    prop_schema,
                    sheet,
                    prefix,
                    current_path
                )
                descriptors.extend(nested)
    
    # Handle array at root level
    elif schema_type == 'array':
        items_schema = schema.get('items', {})
        nested = _flatten_json_schema(
            items_schema,
            sheet,
            prefix,
            f"{path}[]" if path else "[]"
        )
        descriptors.extend(nested)
    
    return descriptors


__all__ = [
    "ExcelSummary",
    "ExcelObjectDescriptor",
    "PresentationSummary",
    "WordField",
    "WordSummary",
    "SwaggerOperationDescriptor",
    "SwaggerSummary",
    "parse_excel",
    "parse_presentation",
    "parse_swagger",
    "parse_word_document",
    "summarize_excel_analysis",
    "summarize_word_analysis",
]
